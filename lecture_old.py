import os
import yt_dlp
import moviepy.editor as mp
import pptx
from summa.summarizer import summarize
import re
from moviepy.editor import AudioFileClip, concatenate_videoclips, concatenate_audioclips,ImageClip, concatenate_videoclips
from TTS.api import TTS
import glob
import zipfile
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from faster_whisper import WhisperModel
import comtypes.client
from pydub import AudioSegment

# Define static path
ZIP_FILE_PATH = "voice_dataset.zip"
EXTRACTED_FOLDER = "voice_dataset"
AUDIO_FOLDER = os.path.join(EXTRACTED_FOLDER, "recording_samples_wav")
PPT_IMAGES_FOLDER = "static/lecture/ppt_images"

# Step 1: Download YouTube Video
def download_youtube_video(youtube_url, output_path="downloaded_video.mp4"):
    ydl_opts = {'format': 'best', 'outtmpl': output_path, 'nocheckcertificate': True}
    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
        ydl.download([youtube_url])
    print(f"✅ YouTube video downloaded: {output_path}")
    return output_path

# Step 2: Extract Audio
def extract_audio(video_path, audio_path=None):
    if audio_path is None:
        audio_path = os.path.join(os.path.dirname(video_path), "audio.wav")
    
    video = mp.VideoFileClip(video_path)
    if video.audio:
        video.audio.write_audiofile(audio_path)
        print(f"✅ Audio extracted successfully to: {audio_path}")
        return audio_path
    print("⚠️ No audio found.")
    return None
    
# Step 3: Transcribe Audio
def transcribe_audio(audio_path):
    model = WhisperModel("tiny", compute_type="int8", device="cpu")  # Use tiny model for speed
    segments, _ = model.transcribe(audio_path)
    transcript = " ".join(segment.text for segment in segments)
    return transcript

# Step 4: Summarize Text
def create_phrased_summary(transcript):
    summary = summarize(transcript, ratio=0.4)
    print("✅ Summary generated.")
    print("Summary:",summary)
    return re.split(r'(?<=[.!?])\s+', summary)

def extract_keywords(sentence):
    return [word for word in sentence.split() if len(word) > 4][:3]

# Step 5: Create PPT from Summary_phrases
def create_ppt(summary_phrases, ppt_path="summary_presentation.pptx"):
    presentation = pptx.Presentation()
    points_per_slide = 3
    for i in range(0, len(summary_phrases), points_per_slide):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        # Slide title
        title = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
        title_frame = title.text_frame
        title_frame.text = f"Key Takeaways - Slide {i//points_per_slide + 1}"
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
       
       # Content
        content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
        content_frame = content.text_frame
        content_frame.word_wrap = True
        
        for point in summary_phrases[i:i+points_per_slide]:
            p = content_frame.add_paragraph()
            keywords = extract_keywords(point)
            
            run = p.add_run()
            run.text = "• "  # Bullet point
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(50, 50, 50)
            
            words = point.split()
            for word in words:
                run = p.add_run()
                run.text = f"{word} "
                run.font.size = Pt(32)
                run.font.name = "Arial"
                run.font.color.rgb = RGBColor(50, 50, 50)
                
                if word in keywords:
                    run.font.color.rgb = RGBColor(0, 128, 0)  # Green color for keywords
            
            p.space_after = Pt(12)
    
    presentation.save(ppt_path)
    return(ppt_path)

# Step 6: Convert PPT Slides to Images
def convert_ppt_to_images(ppt_path):
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1
    ppt = powerpoint.Presentations.Open(os.path.abspath(ppt_path))
    ppt.SaveAs(os.path.abspath(PPT_IMAGES_FOLDER), 17)
    ppt.Close()
    powerpoint.Quit()
    print("✅ PPT slides converted to images.")

# Step 7: Unzip and Select Reference Audio for Cloning
def unzip_dataset():
    if not os.path.exists(EXTRACTED_FOLDER):
        with zipfile.ZipFile(ZIP_FILE_PATH, "r") as zip_ref:
            zip_ref.extractall(EXTRACTED_FOLDER)
        print("✅ Dataset extracted.")

def get_reference_audio():
    audio_files = sorted(glob.glob(os.path.join(AUDIO_FOLDER, "*.wav")))
    return audio_files[0] if audio_files else None

def split_text(text, max_chars=100):
    """Splits text into chunks of max_chars, ensuring words are not cut off."""
    chunks = []
    start = 0

    while start < len(text):
        end = start + max_chars
        
        # If end is not at a space, move back to the last space
        if end < len(text) and text[end] != " ":
            end = text.rfind(" ", start, end)

            # If no space found, just take max_chars to avoid infinite loop
            if end == -1:
                end = start + max_chars
        
        chunks.append(text[start:end].strip())
        start = end

    return chunks

def generate_speech(summary, output_path, reference_voice_sample):
    if isinstance(summary, list):
        summary = " ".join(summary)
    # 1️⃣ Check if summary is valid
    if not summary or not summary.strip():
        print("⚠️ Summary is empty!")
        return None
    print("Summary after: ", summary)
    print("Type of summary:", type(summary))
    # 2️⃣ Check if reference voice sample exists
    if reference_voice_sample is None or not os.path.exists(reference_voice_sample):
        print("⚠️ Reference voice sample is missing or invalid!")
        return None
    # 3️⃣ Split text properly
    text_chunks = split_text(summary, max_chars=100)
    print(f"📝 Generating speech for {len(text_chunks)} chunks.")
    tts = TTS(model_name="tts_models/multilingual/multi-dataset/xtts_v2").to("cpu")
    output_files = []
    print("Count of text chunks:", len(text_chunks))
    print("Text chunks:", text_chunks)
    for i, chunk in enumerate(text_chunks):
        if not chunk.strip():  # Skip empty chunks
            continue
        chunk_path = output_path.replace(".wav", f"_{i}.wav")
        print(f"🔊 Synthesizing chunk {i}: {chunk}")
        # 4️⃣ Ensure valid TTS call
        tts.tts_to_file(text=chunk, file_path=chunk_path, speaker_wav=reference_voice_sample, language="en")
        output_files.append(chunk_path)
    print(f"✅ Speech synthesis complete. Generated {len(output_files)} files.")
    # 5️⃣ Merge all audio files into one
    final_output = merge_audio_chunks(output_files, output_path)
    return final_output  # Return the final merged audio path

# Function to merge all audio chunks into a single .wav file
def merge_audio_chunks(audio_files, final_output_path):
    if not audio_files:
        print("⚠️ No audio chunks found to merge!")
        return None

    combined = AudioSegment.empty()

    for file in audio_files:
        audio = AudioSegment.from_wav(file)
        combined += audio  # Append each audio file

    # Export the final merged file
    combined.export(final_output_path, format="wav")
    print(f"✅ Merged audio saved as: {final_output_path}")

    # Optional: Delete temporary chunk files
    for file in audio_files:
        os.remove(file)

    return final_output_path
  # Return list of generated audio files

# Step 9: Create Video from Images
def create_video_from_images(image_folder, output_video="summary_video.mp4", slide_duration=20):
    image_files = sorted(glob.glob(f"{image_folder}/*.JPG"))
    if not image_files:
        print("⚠️ No images found!")
        return None
    clips = [ImageClip(img).set_duration(slide_duration) for img in image_files]
    final_clip = concatenate_videoclips(clips, method="compose")
    final_clip.write_videofile(output_video, codec="libx264", fps=24)
    print(f"✅ Video created successfully: {output_video}")
    return output_video

# Step 10: Overlay Speech on Video
def overlay_audio_on_video(video_path, audio_path, final_output="final_summary_video.mp4"):
    video = mp.VideoFileClip(video_path)
    audio = AudioFileClip(audio_path)

    if audio.duration < video.duration:
        audio = audio.set_duration(video.duration)
    elif audio.duration > video.duration:
        video = video.set_duration(audio.duration)

    final_clip = video.set_audio(audio)
    final_clip.write_videofile(final_output, codec="libx264", fps=24)
    print(f"✅ Final video with voice-over: {final_output}")
    return final_output

# Main function
def main(youtube_url):
    video_path = download_youtube_video(youtube_url)
    audio_path = extract_audio(video_path)
    if not audio_path:
        return
    transcript = transcribe_audio(audio_path)
    summary = create_phrased_summary(transcript)
    ppt_path = create_ppt(summary)
    convert_ppt_to_images(ppt_path)
    unzip_dataset()
    reference_audio = get_reference_audio()
    if not reference_audio:
        print("⚠️ No reference voice sample found!")
        return
    speech_path = "static/lecture/summary_speech.wav"
    generate_speech(summary, speech_path, reference_audio)
    video_from_images = create_video_from_images(PPT_IMAGES_FOLDER)
    if video_from_images and speech_path:
        overlay_audio_on_video(video_from_images, speech_path)

if __name__ == "__main__":
    youtube_video_url = "https://www.youtube.com/watch?v=LYVwbB2-tnE&t=7s"
    main(youtube_video_url)
