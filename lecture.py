import os
import yt_dlp
import moviepy.editor as mp
import pptx
from summa.summarizer import summarize
import re
from moviepy.editor import AudioFileClip, concatenate_videoclips, concatenate_audioclips, ImageClip
from TTS.api import TTS
import glob
import zipfile
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from faster_whisper import WhisperModel
import comtypes.client
from pydub import AudioSegment
import tempfile
from transformers import pipeline

# Define static path
ZIP_FILE_PATH = "voice_dataset.zip"
EXTRACTED_FOLDER = "voice_dataset"
EXTRACTED_FOLDER_FEMALE = "voice_dataset_female"
AUDIO_FOLDER = os.path.join(EXTRACTED_FOLDER, "recording_samples_wav")
AUDIO_FOLDER_FEMALE = os.path.join(EXTRACTED_FOLDER_FEMALE, "recording_samples_wav")
PPT_IMAGES_FOLDER = "static/lecture/ppt_images"

# Load the BART summarizer
bart_summarizer = pipeline("summarization", model="facebook/bart-large-cnn")

# Ensure directories exist
os.makedirs(PPT_IMAGES_FOLDER, exist_ok=True)
os.makedirs("static/lecture", exist_ok=True)

# Step 1: Download YouTube Video
def download_youtube_video(youtube_url, output_path="downloaded_video.mp4"):
    ydl_opts = {
        'format': 'best',
        'outtmpl': output_path,
        'nocheckcertificate': True,
        'quiet': True
    }
    try:
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            ydl.download([youtube_url])
        print(f"✅ YouTube video downloaded: {output_path}")
        return output_path
    except Exception as e:
        print(f"❌ Error downloading YouTube video: {e}")
        return None

# Step 2: Extract Audio
def extract_audio(video_path, audio_path=None):
    if audio_path is None:
        audio_path = os.path.join(os.path.dirname(video_path), "audio.wav")
    
    try:
        video = mp.VideoFileClip(video_path)
        if video.audio:
            video.audio.write_audiofile(audio_path, verbose=False, logger=None)
            print(f"✅ Audio extracted successfully to: {audio_path}")
            return audio_path
        print("⚠️ No audio found in the video.")
        return None
    except Exception as e:
        print(f"❌ Error extracting audio: {e}")
        return None
    
# Step 3: Transcribe Audio
def transcribe_audio(audio_path):
    try:
        model = WhisperModel("tiny", compute_type="int8", device="cpu")
        segments, _ = model.transcribe(audio_path)
        transcript = " ".join(segment.text for segment in segments)
        return transcript
    except Exception as e:
        print(f"❌ Error transcribing audio: {e}")
        return ""

# Step 4: Summarize Text
def create_phrased_summary(transcript, chunk_size=1024):
    chunks = []
    while len(transcript) > chunk_size:
        split_point = transcript[:chunk_size].rfind('.')
        if split_point == -1:
            split_point = chunk_size
        chunks.append(transcript[:split_point + 1].strip())
        transcript = transcript[split_point + 1:].strip()
    if transcript:
        chunks.append(transcript)

    print(f"🧠 Total chunks to summarize: {len(chunks)}")

    full_summary = ""
    for chunk in chunks:
        summary = bart_summarizer(chunk, max_length=180, min_length=40, do_sample=False)[0]['summary_text']
        full_summary += summary.strip() + " "

    sentences = re.split(r'(?<=[.!?])\s+', full_summary.strip())
    refined_phrases = [re.sub(r'[^\w\s]', '', sentence).strip() for sentence in sentences if len(sentence.strip()) > 0]
    return refined_phrases

def extract_keywords(sentence):
    return [word for word in sentence.split() if len(word) > 4][:3]

# Step 5: Create PPT from Summary_phrases
def create_ppt(summary_phrases, ppt_path="summary_presentation.pptx"):
    try:
        presentation = pptx.Presentation()
        points_per_slide = 3
        for i in range(0, len(summary_phrases), points_per_slide):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            # Slide title
            title = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
            title_frame = title.text_frame
            title_frame.text = f"Key Takeaways - Slide {i//points_per_slide + 1}"
            title_frame.paragraphs[0].font.size = Pt(44)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
           
           # Content
            content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
            content_frame = content.text_frame
            content_frame.word_wrap = True
            
            for point in summary_phrases[i:i+points_per_slide]:
                p = content_frame.add_paragraph()
                keywords = extract_keywords(point)
                
                run = p.add_run()
                run.text = "• "  # Bullet point
                run.font.size = Pt(32)
                run.font.bold = True
                run.font.color.rgb = RGBColor(50, 50, 50)
                
                words = point.split()
                for word in words:
                    run = p.add_run()
                    run.text = f"{word} "
                    run.font.size = Pt(32)
                    run.font.name = "Arial"
                    run.font.color.rgb = RGBColor(50, 50, 50)
                    
                    if word in keywords:
                        run.font.color.rgb = RGBColor(0, 128, 0)  # Green color for keywords
                
                p.space_after = Pt(12)
        
        presentation.save(ppt_path)
        return ppt_path
    except Exception as e:
        print(f"❌ Error creating PPT: {e}")
        return None

# Step 6: Convert PPT Slides to Images
def convert_ppt_to_images(ppt_path):
    try:
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1
        ppt = powerpoint.Presentations.Open(os.path.abspath(ppt_path))
        ppt.SaveAs(os.path.abspath(PPT_IMAGES_FOLDER), 17)  # 17 is for JPG format
        ppt.Close()
        powerpoint.Quit()
        print("✅ PPT slides converted to images.")
        return True
    except Exception as e:
        print(f"❌ Error converting PPT to images: {e}")
        return False

# Step 7: Unzip and Select Reference Audio for Cloning
def unzip_dataset():
    try:
        if not os.path.exists(EXTRACTED_FOLDER):
            with zipfile.ZipFile(ZIP_FILE_PATH, "r") as zip_ref:
                zip_ref.extractall(EXTRACTED_FOLDER)
            print("✅ Dataset extracted.")
        return True
    except Exception as e:
        print(f"❌ Error extracting dataset: {e}")
        return False

def get_reference_audio():
    try:
        audio_files = sorted(glob.glob(os.path.join(AUDIO_FOLDER, "*.wav")))
        return audio_files[0] if audio_files else None
    except Exception as e:
        print(f"❌ Error getting reference audio: {e}")
        return None
    
def get_reference_audio_female():
    try:
        audio_files = sorted(glob.glob(os.path.join(AUDIO_FOLDER_FEMALE, "*.wav")))
        return audio_files[0] if audio_files else None
    except Exception as e:
        print(f"❌ Error getting reference audio: {e}")
        return None

def split_text(text, max_chars=100):
    """Splits text into chunks of max_chars, ensuring words are not cut off."""
    chunks = []
    start = 0

    while start < len(text):
        end = start + max_chars
        
        # If end is not at a space, move back to the last space
        if end < len(text) and text[end] != " ":
            end = text.rfind(" ", start, end)

            # If no space found, just take max_chars to avoid infinite loop
            if end == -1:
                end = start + max_chars
        
        chunks.append(text[start:end].strip())
        start = end

    return chunks

def generate_speech(summary, output_path, reference_voice_sample):
    try:
        if isinstance(summary, list):
            summary = " ".join(summary)
        
        if not summary or not summary.strip():
            print("⚠️ Summary is empty!")
            return None
            
        print("Summary after: ", summary)
        print("Type of summary:", type(summary))
        
        if reference_voice_sample is None or not os.path.exists(reference_voice_sample):
            print("⚠️ Reference voice sample is missing or invalid!")
            return None
            
        text_chunks = split_text(summary, max_chars=100)
        print(f"📝 Generating speech for {len(text_chunks)} chunks.")
        
        tts = TTS(model_name="tts_models/multilingual/multi-dataset/xtts_v2").to("cpu")
        output_files = []
        print("Count of text chunks:", len(text_chunks))
        print("Text chunks:", text_chunks)
        
        # Create a temporary directory for audio chunks
        with tempfile.TemporaryDirectory() as temp_dir:
            for i, chunk in enumerate(text_chunks):
                if not chunk.strip():
                    continue
                chunk_path = os.path.join(temp_dir, f"chunk_{i}.wav")
                print(f"🔊 Synthesizing chunk {i}: {chunk}")
                
                try:
                    tts.tts_to_file(
                        text=chunk,
                        file_path=chunk_path,
                        speaker_wav=reference_voice_sample,
                        language="en"
                    )
                    output_files.append(chunk_path)
                except Exception as e:
                    print(f"⚠️ Error synthesizing chunk {i}: {e}")
                    continue
            
            if not output_files:
                print("⚠️ No audio chunks were generated successfully!")
                return None
                
            # Merge all audio chunks
            final_output = merge_audio_chunks(output_files, output_path)
            return final_output
    except Exception as e:
        print(f"❌ Error in speech generation: {e}")
        return None

def merge_audio_chunks(audio_files, final_output_path):
    try:
        if not audio_files:
            print("⚠️ No audio chunks found to merge!")
            return None

        combined = AudioSegment.empty()

        for file in audio_files:
            try:
                audio = AudioSegment.from_wav(file)
                combined += audio
            except Exception as e:
                print(f"⚠️ Error processing audio chunk {file}: {e}")
                continue

        # Export the final merged file
        combined.export(final_output_path, format="wav")
        print(f"✅ Merged audio saved as: {final_output_path}")
        return final_output_path
    except Exception as e:
        print(f"❌ Error merging audio chunks: {e}")
        return None
  
# Step 9: Create Video from Images
def create_video_from_images(image_folder, output_video="summary_video.mp4", slide_duration=20):
    try:
        image_files = sorted(glob.glob(os.path.join(image_folder, "*.JPG")))
        if not image_files:
            print("⚠️ No images found!")
            return None
            
        clips = [ImageClip(img).set_duration(slide_duration) for img in image_files]
        final_clip = concatenate_videoclips(clips, method="compose")
        
        # Write video without audio first
        temp_video = "temp_video.mp4"
        final_clip.write_videofile(
            temp_video,
            codec="libx264",
            fps=24,
            audio=False,
            verbose=False,
            logger=None
        )
        
        # Rename to final output
        if os.path.exists(output_video):
            os.remove(output_video)
        os.rename(temp_video, output_video)
        
        print(f"✅ Video created successfully: {output_video}")
        return output_video
    except Exception as e:
        print(f"❌ Error creating video from images: {e}")
        return None

# Step 10: Overlay Speech on Video
def overlay_audio_on_video(video_path, audio_path, final_output="final_summary_video.mp4"):
    try:
        # Load video and audio
        video = mp.VideoFileClip(video_path)
        audio = mp.AudioFileClip(audio_path)
        
        # Ensure audio is the same duration as video
        if audio.duration > video.duration:
            # If audio is longer, extend video with last frame
            last_frame = video.get_frame(video.duration - 0.1)
            extension = mp.ImageClip(last_frame, duration=audio.duration - video.duration)
            video = mp.concatenate_videoclips([video, extension])
        elif audio.duration < video.duration:
            # If video is longer, extend audio with silence
            silence = mp.AudioClip(lambda t: 0, duration=video.duration - audio.duration)
            audio = mp.concatenate_audioclips([audio, silence])
        
        # Set audio to video
        final_clip = video.set_audio(audio)
        
        # Write final video
        temp_output = "temp_final.mp4"
        final_clip.write_videofile(
            temp_output,
            codec="libx264",
            fps=24,
            audio_codec="aac",
            verbose=False,
            logger=None
        )
        
        # Rename to final output
        if os.path.exists(final_output):
            os.remove(final_output)
        os.rename(temp_output, final_output)
        
        print(f"✅ Final video with voice-over: {final_output}")
        return final_output
    except Exception as e:
        print(f"❌ Error overlaying audio on video: {e}")
        return None

# Main function
def main(youtube_url):
    try:
        # Step 1: Download YouTube video
        video_path = download_youtube_video(youtube_url)
        if not video_path:
            return
            
        # Step 2: Extract audio
        audio_path = extract_audio(video_path)
        if not audio_path:
            return
            
        # Step 3: Transcribe audio
        transcript = transcribe_audio(audio_path)
        if not transcript:
            print("⚠️ No transcript generated!")
            return
            
        # Step 4: Create summary
        summary = create_phrased_summary(transcript)
        if not summary:
            print("⚠️ No summary generated!")
            return
            
        # Step 5: Create PPT
        ppt_path = create_ppt(summary)
        if not ppt_path:
            return
            
        # Step 6: Convert PPT to images
        if not convert_ppt_to_images(ppt_path):
            return
            
        # Step 7: Unzip dataset and get reference audio
        if not unzip_dataset():
            return
            
        reference_audio = get_reference_audio()
        if not reference_audio:
            print("⚠️ No reference voice sample found!")
            return
            
        # Step 8: Generate speech
        speech_path = "static/lecture/summary_speech.wav"
        if not generate_speech(summary, speech_path, reference_audio):
            return
            
        # Step 9: Create video from images
        video_from_images = create_video_from_images(PPT_IMAGES_FOLDER)
        if not video_from_images:
            return
            
        # Step 10: Overlay audio on video
        if not overlay_audio_on_video(video_from_images, speech_path):
            return
            
        print("🎉 All steps completed successfully!")
    except Exception as e:
        print(f"❌ Unexpected error in main function: {e}")

if __name__ == "__main__":
    youtube_video_url = "https://www.youtube.com/watch?v=LYVwbB2-tnE&t=7s"
    main(youtube_video_url)