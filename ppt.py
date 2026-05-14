import whisper
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from transformers import pipeline
from rouge_score import rouge_scorer
import matplotlib.pyplot as plt
import os
import re

from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader

# Add path to ffmpeg if needed
os.environ["PATH"] += os.pathsep + r"C:\Users\PRIYA\Downloads\ffmpeg"

# Load the BART summarizer
bart_summarizer = pipeline("summarization", model="facebook/bart-large-cnn")

# Function to extract transcript using Whisper
def extract_transcript(video_path):
    model = whisper.load_model("base", device="cpu")
    result = model.transcribe(video_path)
    return result["text"]

# Keyword highlighter for PPT
def extract_keywords(sentence):
    words = sentence.split()
    keywords = [word for word in words if len(word) > 4]
    return keywords[:3]

# Summarization into phrases from transcript
def create_phrased_summary(transcript, chunk_size=1024):
    chunks = []
    while len(transcript) > chunk_size:
        split_point = transcript[:chunk_size].rfind('.')
        if split_point == -1:
            split_point = chunk_size
        chunks.append(transcript[:split_point + 1].strip())
        transcript = transcript[split_point + 1:].strip()
    if transcript:
        chunks.append(transcript)

    print(f"🧠 Total chunks to summarize: {len(chunks)}")

    full_summary = ""
    for chunk in chunks:
        summary = bart_summarizer(chunk, max_length=180, min_length=40, do_sample=False)[0]['summary_text']
        full_summary += summary.strip() + " "

    sentences = re.split(r'(?<=[.!?])\s+', full_summary.strip())
    refined_phrases = [re.sub(r'[^\w\s]', '', sentence).strip() for sentence in sentences if len(sentence.strip()) > 0]
    return refined_phrases

# PPT generator from key points
def create_ppt(summary_phrases, ppt_path):
    presentation = pptx.Presentation()
    points_per_slide = 3
    max_slides = 8
    for i in range(0, min(len(summary_phrases), max_slides * points_per_slide), points_per_slide):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        title = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
        title_frame = title.text_frame
        title_frame.text = f"Key Takeaways - Slide {i // points_per_slide + 1}"
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
        content_frame = content.text_frame
        content_frame.word_wrap = True

        for point in summary_phrases[i:i+points_per_slide]:
            p = content_frame.add_paragraph()
            keywords = extract_keywords(point)

            run = p.add_run()
            run.text = "• "
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(50, 50, 50)

            words = point.split()
            for word in words:
                run = p.add_run()
                run.text = f"{word} "
                run.font.size = Pt(32)
                run.font.name = "Arial"
                run.font.color.rgb = RGBColor(50, 50, 50)

                if word in keywords:
                    run.font.color.rgb = RGBColor(0, 128, 0)

            p.space_after = Pt(12)

    presentation.save(ppt_path)
    print(f"✅ PPT saved to: {ppt_path}")

# Normalize text for ROUGE scoring
def normalize(text):
    text = text.lower()
    text = re.sub(r"[^\w\s]", "", text)
    return text

# Evaluation of summary vs ground truth (transcript)
def evaluate_summary(predicted, reference):
    predicted = normalize(predicted)
    reference = normalize(reference)

    scorer = rouge_scorer.RougeScorer(['rouge1', 'rouge2', 'rougeL'], use_stemmer=True)
    scores = scorer.score(reference, predicted)

    precision = scores['rouge1'].precision
    recall = scores['rouge1'].recall
    f1 = scores['rouge1'].fmeasure

    print("\n📊 Evaluation Metrics:")
    print(f"ROUGE-1 Precision: {precision:.4f}")
    print(f"ROUGE-1 Recall:    {recall:.4f}")
    print(f"ROUGE-1 F1 Score:  {f1:.4f}")
    print(f"ROUGE-2 F1 Score:  {scores['rouge2'].fmeasure:.4f}")
    print(f"ROUGE-L F1 Score:  {scores['rougeL'].fmeasure:.4f}")

    return {
        'ROUGE-1 Precision': precision,
        'ROUGE-1 Recall': recall,
        'ROUGE-1 F1': f1,
        'ROUGE-2 F1': scores['rouge2'].fmeasure,
        'ROUGE-L F1': scores['rougeL'].fmeasure
    }

# Visualization of evaluation metrics
def visualize_metrics(metrics_dict, image_path="metrics_plot.png"):
    labels = list(metrics_dict.keys())
    values = [metrics_dict[k] for k in labels]

    plt.figure(figsize=(10, 5))
    bars = plt.bar(labels, values, color='skyblue')
    for bar in bars:
        yval = bar.get_height()
        plt.text(bar.get_x() + 0.1, yval + 0.01, f"{yval:.2f}")

    plt.title("Summary Evaluation Metrics")
    plt.ylabel("Score")
    plt.ylim(0, 1)
    plt.grid(axis='y', linestyle='--', alpha=0.6)
    plt.tight_layout()
    plt.savefig(image_path)
    plt.show()

# PDF report generator
def generate_pdf_report(metrics, chart_path, pdf_path="ppt_evaluation_metric.pdf"):
    c = canvas.Canvas(pdf_path, pagesize=letter)
    width, height = letter

    c.setFont("Helvetica-Bold", 20)
    c.drawString(50, height - 50, "Summary Evaluation Report")

    c.setFont("Helvetica", 12)
    y_position = height - 100
    for key, value in metrics.items():
        c.drawString(60, y_position, f"{key}: {value:.4f}")
        y_position -= 20

    # Insert the chart
    img = ImageReader(chart_path)
    img_width = 400
    img_height = 250
    c.drawImage(img, 100, 150, width=img_width, height=img_height)

    c.showPage()
    c.save()
    print(f"📄 PDF report saved to: {pdf_path}")

# Main process
if __name__ == "__main__":
    video_path = "downloaded_video.mp4"
    ppt_path = "video_summary.pptx"
    chart_path = "metrics_plot.png"
    pdf_path = "ppt_evaluation_metric.pdf"

    print("🔍 Extracting transcript from video...")
    transcript = extract_transcript(video_path)

    print("✂️ Summarizing transcript using BART...")
    phrases = create_phrased_summary(transcript)

    print("📊 Generating PowerPoint presentation...")
    create_ppt(phrases, ppt_path)

    # Combine all phrases into one string for evaluation
    generated_summary = " ".join(phrases)
    reference_summary = transcript

    print("🧪 Evaluating how well the generated summary captures the transcript...")
    metrics = evaluate_summary(generated_summary, reference_summary)

    print("📈 Visualizing the metrics...")
    visualize_metrics(metrics, chart_path)

    print("📄 Generating PDF evaluation report...")
    generate_pdf_report(metrics, chart_path, pdf_path)

    print(f"🎉 Done! Presentation saved as: {ppt_path}, PDF saved as: {pdf_path}")
